from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_slide(prs, title_text, content_text, image_path=None):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = title_text
    
    # Content
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = content_text
    
    if image_path:
        # Add image to the right side
        left = Inches(5)
        top = Inches(2)
        width = Inches(4.5)
        slide.shapes.add_picture(image_path, left, top, width=width)
        # Adjust text frame width
        body_shape.width = Inches(4.5)

prs = Presentation()

# Slide 1: Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.shapes.placeholders[1]
title.text = "Uçuş Veri Seti Analiz Sunumu"
subtitle.text = "Veri Bilimi Ödevi - Soru 2\nAnaliz: Describe, Outlier, Correlation, Histogram"

# Slide 2: Temel İstatistikler (Describe)
add_slide(prs, "1. Temel İstatistikler (Describe)", 
          "- Ortalama Fiyat: 20,889\n"
          "- Ortalama Süre: 12.22 Saat\n"
          "- Fiyat Standart Sapması: 22,697 (Yüksek varyans)\n"
          "- Minimum Fiyat: 1,105\n"
          "- Maksimum Fiyat: 123,071")

# Slide 3: Aykırı Değer Analizi (Outlier)
add_slide(prs, "2. Aykırı Değerler (IQR Yöntemi)", 
          "- Duration (Süre): 2,110 aykırı değer tespit edildi.\n"
          "- Price (Fiyat): 123 aykırı değer tespit edildi.\n"
          "- Days Left: Aykırı değer saptanmadı.\n\n"
          "Yorum: Fiyattaki aykırı değerler muhtemelen business class veya son dakika biletlerinden kaynaklanmaktadır.")

# Slide 4: Korelasyon Analizi
add_slide(prs, "3. Değişkenler Arası Korelasyon", 
          "- Süre & Fiyat: 0.204 (Düşük Pozitif)\n"
          "- Kalan Gün & Fiyat: -0.091 (Çok Düşük Negatif)\n\n"
          "Bulgu: Uçuş süresi arttıkça fiyatın artma eğilimi vardır ancak ilişki zayıftır. Kalan gün sayısı bu örneklemde fiyatı tek başına güçlü etkilememektedir.")

# Slide 5: Veri Dağılımı (Histogramlar)
add_slide(prs, "4. Veri Dağılımı ve Görselleştirme", 
          "- Duration: Sağdan çarpık dağılım.\n"
          "- Price: Bimodal (iki tepeli) dağılım eğilimi.\n"
          "- Days Left: Homojen (üniform) dağılım.", 
          "histograms.png")

prs.save('Analiz_Sunumu.pptx')
print("Sunum 'Analiz_Sunumu.pptx' adıyla başarıyla oluşturuldu.")
