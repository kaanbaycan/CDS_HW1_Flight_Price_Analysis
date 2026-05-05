from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation('Gelismis_Analiz_Sunumu.pptx')

# Add a specific slide for Detailed IQR
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Detaylı IQR Hesaplamaları"

content = (
    "Aykırı değer tespiti için kullanılan sınırlar:\n\n"
    "1. Duration (Süre):\n"
    "   - IQR: 9.34 | Üst Sınır: 30.18 | Aykırı: 2,110\n"
    "2. Days Left (Kalan Gün):\n"
    "   - IQR: 23.0 | Üst Sınır: 72.50 | Aykırı: 0\n"
    "3. Price (Fiyat):\n"
    "   - IQR: 37,738 | Üst Sınır: 99,128 | Aykırı: 123\n\n"
    "Not: Alt sınırlar negatif çıktığı için (fiyat/süre negatif olamaz) sadece üst sınırlar anlamlıdır."
)

body_shape = slide.placeholders[1]
body_shape.text_frame.text = content

prs.save('Gelismis_Analiz_Sunumu.pptx')
print("Sunum detaylı IQR bilgileriyle güncellendi.")
