from pptx import Presentation
from pptx.util import Inches, Pt

def add_image_slide(prs, title_text, content_text, image_path):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    
    body_shape = slide.placeholders[1]
    body_shape.text_frame.text = content_text
    body_shape.width = Inches(4.5)
    
    left = Inches(5)
    top = Inches(1.5)
    width = Inches(4.5)
    slide.shapes.add_picture(image_path, left, top, width=width)

prs = Presentation()

# Slide 1: Title
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
title_slide.shapes.title.text = "Gelişmiş Veri Analizi ve Kümeleme"
title_slide.placeholders[1].text = "Uçuş Segmentasyonu ve Varyans Analizi"

# Slide 2: Segmentasyon (Clustering)
add_image_slide(prs, "K-Means Kümeleme Analizi", 
                "Fiyat ve Süre bazlı 4 ana segment oluşturuldu:\n\n"
                "- Cluster 3: En düşük varyans (Std: 3,777)\n"
                "- Kümeleme sonrası grup içi varyans %70'e varan oranlarda azaltıldı.\n"
                "- Bu sayede daha homojen fiyat grupları elde edildi.", 
                "cluster_scatter.png")

# Slide 3: Dağılım Analizi
add_image_slide(prs, "Grup Bazlı Fiyat Dağılımı", 
                "Yoğunluk (KDE) Grafiği Bulguları:\n\n"
                "- Her küme kendine özgü bir fiyat aralığını temsil eder.\n"
                "- Kümelerin ayrışması, segmentasyonun başarısını doğrulamaktadır.", 
                "cluster_kde.png")

# Slide 4: Sınıf Bazlı Analiz
add_image_slide(prs, "Kabin Sınıfı Etkisi", 
                "Fiyat üzerindeki en büyük belirleyici:\n\n"
                "- Ekonomi ve Business sınıfları arasındaki keskin ayrım.\n"
                "- Business sınıfı kendi içinde daha yüksek varyansa sahiptir.", 
                "class_boxplot.png")

prs.save('Gelismis_Analiz_Sunumu.pptx')
print("Gelişmiş sunum oluşturuldu.")
